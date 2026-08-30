#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""文档转换引擎：五级引擎链（零 Qt / 零 FastAPI 依赖）。

从旧 doc_converter.py 的 DocConversionWorker 剥离：
  - 保留五级降级链：原生 Python 库 → pdf2docx → python-pptx → pandoc → HTML→PDF；
  - 保留引擎探测函数：_probe_native_engines / _find_pandoc / _find_wkhtmltopdf；
  - 原 _run_subprocess_with_cancel（依赖 Qt 的 TaskControlState 与 get_settings_int）
    替换为独立的 _run_subprocess(cmd, should_abort_cb)，取消通过轮询回调实现；
  - convert_doc(task, progress_cb, should_abort_cb) 与旧 _convert_single 行为一致，
    返回 (bool, message)，取消时消息含「用户取消」并删除半成品。
"""

import os
import shutil
import subprocess
from collections import deque
from dataclasses import dataclass
from typing import Callable, Optional, Tuple

# =====================================================================
# 格式预设表（沿用旧 doc_converter.py）
# =====================================================================

SUPPORTED_DOC_FORMATS = {
    "pdf":     {"ext": ".pdf",  "desc": "PDF — 通用电子文档（跨平台首选）"},
    "docx":    {"ext": ".docx", "desc": "DOCX — Microsoft Word 现代格式"},
    "doc":     {"ext": ".doc",  "desc": "DOC — 旧版 Word（建议转 DOCX）"},
    "odt":     {"ext": ".odt",  "desc": "ODT — 开放文档文本文档"},
    "xlsx":    {"ext": ".xlsx", "desc": "XLSX — Microsoft Excel 现代表格"},
    "xls":     {"ext": ".xls",  "desc": "XLS — 旧版 Excel（建议转 XLSX）"},
    "csv":     {"ext": ".csv",  "desc": "CSV — 逗号分隔值表格"},
    "pptx":    {"ext": ".pptx", "desc": "PPTX — Microsoft PowerPoint 现代演示"},
    "ppt":     {"ext": ".ppt",  "desc": "PPT — 旧版 PowerPoint"},
    "txt":     {"ext": ".txt",  "desc": "TXT — 纯文本"},
    "md":      {"ext": ".md",   "desc": "Markdown — 轻量标记语言"},
    "html":    {"ext": ".html", "desc": "HTML — 网页格式"},
    "rtf":     {"ext": ".rtf",  "desc": "RTF — 富文本格式"},
    "epub":    {"ext": ".epub", "desc": "EPUB — 电子书标准格式"},
}

EXTRA_INPUT_EXTS = [
    "jpg", "jpeg", "jfif", "png", "webp", "bmp", "gif", "tif", "tiff",
    "heic", "heif", "svg",
]

ALL_INPUT_EXTS = set(
    list(SUPPORTED_DOC_FORMATS.keys()) + EXTRA_INPUT_EXTS +
    [e.lstrip(".") for e in [
        ".tex", ".json", ".xml", ".yaml", ".yml", ".ini", ".log", ".py",
        ".js", ".ts", ".c", ".cpp", ".h", ".mdx", ".odp", ".ods", ".epub",
        ".sxc", ".wri", ".wpd", ".pages", ".numbers",
    ]]
)


# =====================================================================
# 数据类 / 工具
# =====================================================================

@dataclass
class DocConversionTask:
    input_path: str
    output_path: str
    output_format: str
    keep_original: bool = True
    pdf_dpi: int = 200
    overwrite: bool = False
    task_id: Optional[int] = None


def _is_image_input(path: str) -> bool:
    ext = os.path.splitext(path)[1].lower().lstrip(".")
    return ext in EXTRA_INPUT_EXTS


def _ext_aware_path(path: str, new_ext: str) -> str:
    return os.path.splitext(path)[0] + new_ext


def _delete_partial(path: str) -> None:
    """取消/失败后删除半成品输出文件。"""
    try:
        if path and os.path.exists(path):
            os.remove(path)
    except Exception:
        pass


# =====================================================================
# 引擎探测（沿用旧 doc_converter.py）
# =====================================================================

def _find_pandoc() -> Optional[str]:
    try:
        return shutil.which("pandoc")
    except Exception:
        return None


def _find_wkhtmltopdf() -> Optional[str]:
    """探测 wkhtmltopdf 二进制（HTML→PDF 专用）。"""
    for name in ("wkhtmltopdf", "wkhtmltopdf.exe"):
        try:
            p = shutil.which(name)
            if p:
                return p
        except Exception:
            pass
    for hard in (
        "/usr/bin/wkhtmltopdf",
        "/usr/local/bin/wkhtmltopdf",
        "/opt/homebrew/bin/wkhtmltopdf",
        r"C:\Program Files\wkhtmltopdf\bin\wkhtmltopdf.exe",
    ):
        if os.path.isfile(hard):
            return hard
    return None


def _probe_native_engines() -> Tuple[dict, dict]:
    """探测当前 Python 环境内的文档处理库可用性。

    返回 (flags, errors)：
      flags[k]   -> bool 是否可用
      errors[k]  -> None 或错误说明，区分「未安装」与「已安装但导入失败」
    """
    import importlib

    specs = {
        "docx2txt": "docx2txt",
        "python_docx": ("docx", "Document"),
        "openpyxl": "openpyxl",
        "fitz": "fitz",
        "pypdf": "pypdf",
        "PIL": ("PIL", "Image"),
        "striprtf": "striprtf",
        "markdown": "markdown",
        "csv": "csv",
        "python_pptx": "pptx",
        "pdf2docx": ("pdf2docx", "parse"),
        "weasyprint": ("weasyprint", "HTML"),
        "pypandoc": "pypandoc",
    }
    flags, errors = {}, {}
    for key, spec in specs.items():
        try:
            if isinstance(spec, tuple):
                module_name, attr = spec
                getattr(importlib.import_module(module_name), attr)
            else:
                importlib.import_module(spec)
            flags[key] = True
            errors[key] = None
        except ImportError as e:
            flags[key] = False
            errors[key] = f"未安装: {getattr(e, 'name', None) or e}"
        except Exception as e:
            flags[key] = False
            errors[key] = f"导入失败({type(e).__name__}): {e}"
    return flags, errors


def probe() -> dict:
    """返回文档引擎探测结果（供 /api/health 与 /api/formats）。"""
    flags, errors = _probe_native_engines()
    return {
        "native_flags": flags,
        "native_errors": errors,
        "pandoc": _find_pandoc(),
        "wkhtmltopdf": _find_wkhtmltopdf(),
    }


# =====================================================================
# 文档引擎
# =====================================================================

class DocEngine:
    """五级引擎链实现（替代旧 DocConversionWorker 的同步部分）。"""

    def __init__(self, native: Optional[dict] = None,
                 pandoc_path: Optional[str] = None,
                 wkhtmltopdf_path: Optional[str] = None):
        flags, _ = _probe_native_engines()
        self.native = native if native is not None else flags
        self.pandoc_path = pandoc_path if pandoc_path is not None else _find_pandoc()
        self.wkhtmltopdf_path = wkhtmltopdf_path if wkhtmltopdf_path is not None else _find_wkhtmltopdf()

    # ---- 子进程执行器（带取消，替代原 _run_subprocess_with_cancel） ----
    def _run_subprocess(self, cmd, should_abort_cb: Callable[[], bool]):
        """带取消能力的子进程执行：轮询等待，中止时 terminate。"""
        try:
            process = subprocess.Popen(
                cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE,
                universal_newlines=True, encoding="utf-8", errors="replace",
            )
        except Exception as e:
            return None, f"启动子进程失败: {e}"
        try:
            while True:
                if should_abort_cb():
                    _terminate(process)
                    return None, "用户取消"
                try:
                    _, stderr = process.communicate(timeout=0.5)
                    break
                except subprocess.TimeoutExpired:
                    if should_abort_cb():
                        _terminate(process)
                        return None, "用户取消"
            if process.returncode == 0:
                return stderr, None
            tail = "\n".join(deque((stderr or "").splitlines(), maxlen=8)) or "(无输出)"
            return stderr, f"exit={process.returncode}\n--- stderr ---\n{tail}"
        except Exception as e:
            return None, f"子进程异常: {e}"

    # ----------------------------------------------------------
    # 1) 原生 Python 引擎（格式特定，速度最快）
    # ----------------------------------------------------------
    def _try_native_python(self, task, should_abort_cb: Callable[[], bool]) -> Tuple[bool, str]:
        src, dst, fmt = task.input_path, task.output_path, task.output_format.lower()
        src_ext = os.path.splitext(src)[1].lower().lstrip(".")

        pure_text_srcs = {"txt", "md", "markdown", "html", "htm", "rtf", "csv", "tsv",
                          "json", "xml", "yaml", "yml", "ini", "log", "py", "js", "ts",
                          "c", "cpp", "h", "mdx", "yaml"}
        pure_text_dsts = {"txt", "md", "html", "htm", "csv", "tsv"}
        if (src_ext in pure_text_srcs or fmt in pure_text_dsts) and not (src_ext.startswith("rtf") and fmt not in {"txt"}):
            # RTF → 文本，走 striprtf
            if src_ext == "rtf" and fmt == "txt" and self.native.get("striprtf"):
                try:
                    from striprtf.striprtf import rtf_to_text
                    with open(src, "r", encoding="utf-8", errors="replace") as f:
                        content = f.read()
                    text = rtf_to_text(content)
                    with open(dst, "w", encoding="utf-8") as fo:
                        fo.write(text)
                    return True, "RTF → TXT (striprtf)"
                except Exception as e:
                    return False, f"striprtf 失败: {e}"
            # Markdown → HTML
            if src_ext in {"md", "markdown"} and fmt == "html" and self.native.get("markdown"):
                try:
                    import markdown as md_lib
                    with open(src, "r", encoding="utf-8", errors="replace") as f:
                        html_body = md_lib.markdown(f.read(), extensions=["fenced_code", "tables"])
                    page = (
                        "<!DOCTYPE html><html><head><meta charset='utf-8'></head>"
                        f"<body style='font-family: sans-serif; line-height: 1.6;'>{html_body}</body></html>"
                    )
                    with open(dst, "w", encoding="utf-8") as fo:
                        fo.write(page)
                    return True, "Markdown → HTML (python-markdown)"
                except Exception as e:
                    return False, f"python-markdown 失败: {e}"
            # CSV → XLSX （若 openpyxl 存在）
            if src_ext == "csv" and fmt == "xlsx" and self.native.get("openpyxl"):
                try:
                    import csv
                    from openpyxl import Workbook
                    wb = Workbook()
                    ws = wb.active
                    ws.title = "Sheet1"
                    with open(src, "r", encoding="utf-8-sig", errors="replace", newline="") as f:
                        reader = csv.reader(f)
                        for i, row in enumerate(reader, start=1):
                            for j, cell in enumerate(row, start=1):
                                ws.cell(row=i, column=j, value=cell)
                    save_dir = os.path.dirname(dst)
                    if save_dir:
                        os.makedirs(save_dir, exist_ok=True)
                    wb.save(dst)
                    return True, "CSV → XLSX (openpyxl)"
                except Exception as e:
                    return False, f"openpyxl 失败: {e}"
            # XLSX → CSV （openpyxl）
            if src_ext == "xlsx" and fmt == "csv" and self.native.get("openpyxl"):
                try:
                    import csv
                    from openpyxl import load_workbook
                    wb = load_workbook(src, read_only=True, data_only=True)
                    ws = wb[wb.sheetnames[0]]
                    save_dir = os.path.dirname(dst)
                    if save_dir:
                        os.makedirs(save_dir, exist_ok=True)
                    with open(dst, "w", encoding="utf-8-sig", newline="") as fo:
                        writer = csv.writer(fo)
                        for row in ws.iter_rows(values_only=True):
                            writer.writerow(["" if c is None else str(c) for c in row])
                    return True, "XLSX → CSV (openpyxl)"
                except Exception as e:
                    return False, f"openpyxl 失败: {e}"
            # DOCX → TXT （docx2txt 优先，否则 python-docx 兜底）
            if src_ext == "docx" and fmt == "txt":
                try:
                    if self.native.get("docx2txt"):
                        import docx2txt
                        text = docx2txt.process(src) or ""
                        with open(dst, "w", encoding="utf-8") as fo:
                            fo.write(text)
                        return True, "DOCX → TXT (docx2txt)"
                    if self.native.get("python_docx"):
                        from docx import Document
                        doc = Document(src)
                        paragraphs = [p.text for p in doc.paragraphs]
                        with open(dst, "w", encoding="utf-8") as fo:
                            fo.write("\n".join(paragraphs))
                        return True, "DOCX → TXT (python-docx)"
                except Exception as e:
                    return False, f"docx 读取失败: {e}"
            # 普通纯文本 → 纯文本（复制/编码转换）
            if (src_ext in pure_text_srcs) and (fmt in pure_text_dsts):
                try:
                    with open(src, "rb") as f:
                        raw = f.read()
                    text = raw.decode("utf-8", errors="replace")
                    save_dir = os.path.dirname(dst)
                    if save_dir:
                        os.makedirs(save_dir, exist_ok=True)
                    with open(dst, "w", encoding="utf-8") as fo:
                        fo.write(text)
                    return True, "文本复制 + UTF-8 标准化"
                except Exception as e:
                    return False, f"文本读写失败: {e}"

        # --- 图片 → PDF（Pillow） ---
        if _is_image_input(src) and fmt == "pdf" and self.native.get("PIL"):
            try:
                from PIL import Image
                with Image.open(src) as img:
                    img.load()
                    if img.mode in ("RGBA", "P", "LA"):
                        img = img.convert("RGB")
                    save_dir = os.path.dirname(dst)
                    if save_dir:
                        os.makedirs(save_dir, exist_ok=True)
                    img.save(dst, "PDF", resolution=task.pdf_dpi)
                    return True, f"图片 → PDF (Pillow {img.size[0]}x{img.size[1]} DPI={task.pdf_dpi})"
            except Exception as e:
                return False, f"Pillow 图片→PDF 失败: {e}"

        # --- PDF → 图片（PyMuPDF） ---
        if src_ext == "pdf" and _is_image_input(dst):
            if self.native.get("fitz"):
                try:
                    import fitz
                    doc = fitz.open(src)
                    save_dir = os.path.dirname(dst)
                    if save_dir:
                        os.makedirs(save_dir, exist_ok=True)
                    mat = fitz.Matrix(task.pdf_dpi / 72.0, task.pdf_dpi / 72.0)
                    base, ext = os.path.splitext(dst)
                    if doc.page_count == 1:
                        page = doc.load_page(0)
                        pix = page.get_pixmap(matrix=mat, alpha=False)
                        pix.save(dst)
                    else:
                        for i in range(doc.page_count):
                            p = doc.load_page(i)
                            pix = p.get_pixmap(matrix=mat, alpha=False)
                            pix.save(f"{base}_p{i+1}{ext}")
                    pages = doc.page_count
                    doc.close()
                    return True, f"PDF → 图片 (PyMuPDF DPI={task.pdf_dpi}, pages={pages})"
                except Exception as e:
                    return False, f"PyMuPDF 失败: {e}"

        # --- PDF 文本提取 (pypdf / fitz) ---
        if src_ext == "pdf" and fmt == "txt":
            if self.native.get("fitz"):
                try:
                    import fitz
                    doc = fitz.open(src)
                    texts = []
                    for i in range(doc.page_count):
                        texts.append(doc.load_page(i).get_text("text") or "")
                    doc.close()
                    with open(dst, "w", encoding="utf-8") as fo:
                        fo.write("\n\n".join(texts))
                    return True, "PDF → TXT (PyMuPDF)"
                except Exception:
                    pass
            if self.native.get("pypdf"):
                try:
                    from pypdf import PdfReader
                    reader = PdfReader(src)
                    with open(dst, "w", encoding="utf-8") as fo:
                        for page in reader.pages:
                            try:
                                fo.write(page.extract_text() or "")
                            except Exception:
                                pass
                            fo.write("\n\n")
                    return True, "PDF → TXT (pypdf)"
                except Exception as e:
                    return False, f"pypdf 失败: {e}"
        return False, "当前格式组合无原生引擎，尝试外部工具…"

    # ----------------------------------------------------------
    # 2) pdf2docx — PDF → DOCX
    # ----------------------------------------------------------
    def _try_pdf2docx(self, task, should_abort_cb: Callable[[], bool]) -> Tuple[bool, str]:
        if not self.native.get("pdf2docx"):
            return False, "未安装 pdf2docx"
        src, dst, fmt = task.input_path, task.output_path, task.output_format.lower()
        src_ext = os.path.splitext(src)[1].lower().lstrip(".")
        if src_ext != "pdf" or fmt != "docx":
            return False, "仅支持 PDF → DOCX"
        save_dir = os.path.dirname(dst)
        if save_dir:
            os.makedirs(save_dir, exist_ok=True)
        try:
            from pdf2docx import parse
            parse(src, dst, start=0, end=None, multi_processing=False)
            if os.path.isfile(dst) and os.path.getsize(dst) > 0:
                return True, "PDF → DOCX (pdf2docx)"
            return False, "pdf2docx 输出文件为空"
        except Exception as e:
            return False, f"pdf2docx 失败: {e}"

    # ----------------------------------------------------------
    # 3) python-pptx — PPTX 提取文本
    # ----------------------------------------------------------
    def _try_python_pptx(self, task, should_abort_cb: Callable[[], bool]) -> Tuple[bool, str]:
        if not self.native.get("python_pptx"):
            return False, "未安装 python-pptx"
        src, dst, fmt = task.input_path, task.output_path, task.output_format.lower()
        src_ext = os.path.splitext(src)[1].lower().lstrip(".")
        if src_ext not in ("pptx", "ppt") or fmt not in ("txt", "pdf"):
            return False, "仅支持 PPTX → TXT / PDF"
        save_dir = os.path.dirname(dst)
        if save_dir:
            os.makedirs(save_dir, exist_ok=True)
        try:
            from pptx import Presentation

            if src_ext == "ppt":
                return False, "python-pptx 不支持 .ppt 老格式"

            prs = Presentation(src)
            if fmt == "txt":
                texts = []
                for i, slide in enumerate(prs.slides, start=1):
                    texts.append(f"===== Slide {i} =====")
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            for para in shape.text_frame.paragraphs:
                                line = "".join(run.text for run in para.runs)
                                if line:
                                    texts.append(line)
                with open(dst, "w", encoding="utf-8") as fo:
                    fo.write("\n".join(texts))
                return True, f"PPTX → TXT (python-pptx, {len(prs.slides)} 页)"
            if fmt == "pdf":
                return False, "python-pptx 无原生 PDF 渲染能力，暂不支持 PPTX→PDF"
        except Exception as e:
            return False, f"python-pptx 失败: {e}"
        return False, "python-pptx 失败"

    # ----------------------------------------------------------
    # 4) HTML → PDF（weasyprint 或 wkhtmltopdf 二选一）
    # ----------------------------------------------------------
    def _try_html_to_pdf(self, task, should_abort_cb: Callable[[], bool]) -> Tuple[bool, str]:
        src, dst, fmt = task.input_path, task.output_path, task.output_format.lower()
        src_ext = os.path.splitext(src)[1].lower().lstrip(".")
        if src_ext not in ("html", "htm") or fmt != "pdf":
            return False, "仅支持 HTML → PDF"
        save_dir = os.path.dirname(dst)
        if save_dir:
            os.makedirs(save_dir, exist_ok=True)

        # 优先 weasyprint（纯 Python，无外部依赖）
        if self.native.get("weasyprint"):
            try:
                from weasyprint import HTML
                HTML(filename=src).write_pdf(dst)
                if os.path.isfile(dst) and os.path.getsize(dst) > 0:
                    return True, "HTML → PDF (weasyprint)"
            except Exception:
                pass  # weasyprint 失败时降级到 wkhtmltopdf

        # 次选 wkhtmltopdf（外部二进制，需系统安装）
        if self.wkhtmltopdf_path:
            try:
                cmd = [
                    self.wkhtmltopdf_path,
                    "--quiet",
                    "--encoding", "utf-8",
                    "--enable-local-file-access",
                    src, dst,
                ]
                stderr, err = self._run_subprocess(cmd, should_abort_cb)
                if err == "用户取消":
                    return False, "用户取消"
                if err is None and os.path.isfile(dst) and os.path.getsize(dst) > 0:
                    return True, "HTML → PDF (wkhtmltopdf)"
                tail = "\n".join(deque((stderr or "").splitlines(), maxlen=6)) or "(无输出)"
                return False, f"wkhtmltopdf 失败\n{err or tail}"
            except Exception as e:
                return False, f"wkhtmltopdf 异常: {e}"

        return False, "未检测到 weasyprint 或 wkhtmltopdf"

    # ----------------------------------------------------------
    # 5) Pandoc 主力引擎（通用文档互转）
    # ----------------------------------------------------------
    def _try_pandoc(self, task, should_abort_cb: Callable[[], bool]) -> Tuple[bool, str]:
        if not self.pandoc_path:
            return False, "未检测到 pandoc"
        src, dst, fmt = task.input_path, task.output_path, task.output_format.lower()
        src_ext = os.path.splitext(src)[1].lower().lstrip(".")
        save_dir = os.path.dirname(dst)
        if save_dir:
            os.makedirs(save_dir, exist_ok=True)
        pandoc_inputs = {
            "md", "markdown", "html", "htm", "docx", "odt", "epub",
            "rst", "latex", "tex", "rtf", "txt", "json", "mediawiki",
            "t2t", "org", "wiki",
        }
        pandoc_outputs = {
            "md", "markdown", "html", "htm", "docx", "odt", "epub",
            "rst", "latex", "tex", "rtf", "txt", "json", "mediawiki",
            "pdf",  # 需配合 LaTeX 引擎（如 tectonic / xelatex）
        }
        if src_ext not in pandoc_inputs or fmt not in pandoc_outputs:
            return False, "格式不在 pandoc 支持范围"
        cmd = [self.pandoc_path, src, "-o", dst, "--standalone"]
        # ── PDF 输出：pandoc 需要配合 LaTeX 引擎 ──
        if fmt == "pdf":
            tectonic = shutil.which("tectonic")
            if tectonic:
                cmd.extend(["--pdf-engine=tectonic"])
            else:
                xelatex = shutil.which("xelatex")
                if xelatex:
                    cmd.extend(["--pdf-engine=xelatex"])
                else:
                    return False, "pandoc → PDF 需安装 tectonic 或 xelatex"
        # ── 中文字体处理（DOCX/HTML/EPUB 输出时） ──
        if fmt in ("docx", "html", "epub"):
            cmd.append("--metadata=lang:zh-CN")
        stderr, err = self._run_subprocess(cmd, should_abort_cb)
        if err == "用户取消":
            return False, "用户取消"
        if err is None:
            if os.path.isfile(dst) and os.path.getsize(dst) > 0:
                return True, f"pandoc → {fmt.upper()}"
            return False, "pandoc 未生成输出文件"
        return False, f"pandoc 失败\n{err}"

    # ----------------------------------------------------------
    # 主入口：五级引擎链
    # ----------------------------------------------------------
    def convert(self, task, progress_cb: Callable[[int], None],
                should_abort_cb: Callable[[], bool]) -> Tuple[bool, str]:
        progress_cb(5)

        save_dir = os.path.dirname(task.output_path)
        if save_dir and not os.path.isdir(save_dir):
            os.makedirs(save_dir, exist_ok=True)

        used = []
        res_msg = ""
        success = False
        last_msg = ""

        def _push(label, s, m):
            nonlocal success, res_msg, last_msg
            used.append(label)
            last_msg = m
            if s:
                success = True
                res_msg = m

        # 1) 原生 Python 引擎
        s, m = self._try_native_python(task, should_abort_cb)
        _push("原生Python", s, m)
        progress_cb(40)

        # 2) PDF → DOCX 专用：pdf2docx
        if not success and not should_abort_cb():
            s, m = self._try_pdf2docx(task, should_abort_cb)
            _push("pdf2docx", s, m)
        progress_cb(55)

        # 3) PPTX 专用：python-pptx
        if not success and not should_abort_cb():
            s, m = self._try_python_pptx(task, should_abort_cb)
            _push("python-pptx", s, m)
        progress_cb(65)

        # 4) pandoc 主力（通用互转）
        if not success and not should_abort_cb():
            s, m = self._try_pandoc(task, should_abort_cb)
            _push("pandoc", s, m)
        progress_cb(80)

        # 5) HTML → PDF 专用（weasyprint / wkhtmltopdf）
        if not success and not should_abort_cb():
            s, m = self._try_html_to_pdf(task, should_abort_cb)
            _push("HTML→PDF", s, m)
        progress_cb(100)

        if m == "用户取消" or should_abort_cb():
            _delete_partial(task.output_path)
            return False, "用户取消"

        msg = res_msg if success else (last_msg or "全部引擎失败")
        if success and used:
            msg = f"{msg}（尝试: {', '.join(used)}）"
        return success, msg


def _terminate(process) -> None:
    """终止子进程（POSIX 先解除暂停再 terminate）。"""
    try:
        if os.name == "posix":
            process.send_signal(subprocess.signal.SIGCONT)
        process.terminate()
    except Exception:
        pass


def convert_doc(task, progress_cb: Callable[[int], None],
                should_abort_cb: Callable[[], bool]) -> Tuple[bool, str]:
    """文档转换入口：创建引擎实例并执行五级引擎链。

    :param task: 鸭子类型任务对象，需含 input_path / output_path / output_format，
                 另有 pdf_dpi / keep_original。
    :param progress_cb: progress_cb(percent:int)
    :param should_abort_cb: should_abort_cb() -> bool
    :return: (成功与否, 消息)
    """
    return DocEngine().convert(task, progress_cb, should_abort_cb)


# =====================================================================
# 格式信息（供 REST /api/formats）
# =====================================================================

def formats_info() -> dict:
    """返回文档模块的格式与预设信息。"""
    outputs = []
    for key, info in SUPPORTED_DOC_FORMATS.items():
        outputs.append({"key": key, "desc": info["desc"], "ext": info["ext"]})
    return {
        "inputs": sorted(ALL_INPUT_EXTS),
        "outputs": outputs,
        "presets": {
            "pdf_dpi": [
                {"label": "150 DPI", "value": 150},
                {"label": "200 DPI（默认）", "value": 200},
                {"label": "300 DPI", "value": 300},
            ],
        },
    }
