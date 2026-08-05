
#!/usr/bin/env python3
# -*- coding: utf-8 -*-
import os
import sys
import shutil
import subprocess
import tempfile
from dataclasses import dataclass
from typing import Optional, List, Tuple
from collections import deque

from PySide6.QtWidgets import (
    QWidget, QVBoxLayout, QHBoxLayout, QLabel, QPushButton, QComboBox,
    QLineEdit, QListWidget, QListWidgetItem,
    QGroupBox, QMessageBox, QMenu, QFileDialog,
    QSizePolicy, QSpinBox, QApplication
)
from PySide6.QtCore import Qt, QThread, Signal, QStandardPaths, QRectF, QSize, QObject, QTimer, Slot
from PySide6.QtGui import (
    QDragEnterEvent, QDropEvent, QPixmap, QPainter, QColor, QPen, QBrush,
    QPainterPath, QIcon
)
from utils import COLORS, BaseConversionWorker, get_cjk_font_qss, ThemeManager


# 文档格式定义
# ============================================================
SUPPORTED_DOC_FORMATS = {
    "pdf":     {"ext": ".pdf",  "desc": "PDF — 通用电子文档（跨平台首选）"},
    "docx":    {"ext": ".docx", "desc": "DOCX — Microsoft Word 现代格式"},
    "doc":     {"ext": ".doc",  "desc": "DOC — 旧版 Word（建议转 DOCX）"},
    "odt":     {"ext": ".odt",  "desc": "ODT — 开放文档文本文档"},
    "xlsx":    {"ext": ".xlsx", "desc": "XLSX — Microsoft Excel 现代表格"},
    "xls":     {"ext": ".xls",  "desc": "XLS — 旧版 Excel（建议转 XLSX）"},
    "csv":     {"ext": ".csv",  "desc": "CSV — 逗号分隔值表格"},
    "pptx":    {"ext": ".pptx", "desc": "PPTX — Microsoft PowerPoint 现代演示"},
    "ppt":     {"ext": ".ppt",  "desc": "PPT — 旧版 PowerPoint"},
    "txt":     {"ext": ".txt",  "desc": "TXT — 纯文本"},
    "md":      {"ext": ".md",   "desc": "Markdown — 轻量标记语言"},
    "html":    {"ext": ".html", "desc": "HTML — 网页格式"},
    "rtf":     {"ext": ".rtf",  "desc": "RTF — 富文本格式"},
    "epub":    {"ext": ".epub", "desc": "EPUB — 电子书标准格式"},
}

# 输入额外支持（图片 → PDF / PDF → 图片 作为常见需求）
EXTRA_INPUT_EXTS = [
    "jpg", "jpeg", "jfif", "png", "webp", "bmp", "gif", "tif", "tiff",
    "heic", "heif", "svg"
]

ALL_INPUT_EXTS = set(
    list(SUPPORTED_DOC_FORMATS.keys()) + EXTRA_INPUT_EXTS +
    [e.lstrip(".") for e in [
        ".tex", ".json", ".xml", ".yaml", ".yml", ".ini", ".log", ".py",
        ".js", ".ts", ".c", ".cpp", ".h", ".mdx", ".odp", ".ods", ".epub",
        ".sxc", ".wri", ".wpd", ".pages", ".numbers"
    ]]
)


# ============================================================
# 数据类
# ============================================================
@dataclass
class DocConversionTask:
    input_path: str
    output_path: str
    output_format: str
    keep_original: bool = True
    pdf_dpi: int = 200
    overwrite: bool = False


@dataclass
class DocConversionResult:
    success: bool
    task: DocConversionTask
    message: str


# ============================================================
# 通用 helper
# ============================================================
def _make_checkbox_icon_pixmap(checked: bool, accent_hex: str, size=22, radius=5):
    pm = QPixmap(size, size)
    pm.fill(Qt.GlobalColor.transparent)
    p = QPainter(pm)
    p.setRenderHint(QPainter.RenderHint.Antialiasing, True)
    if checked:
        p.setBrush(QBrush(QColor(accent_hex)))
        p.setPen(Qt.PenStyle.NoPen)
        path = QPainterPath()
        path.addRoundedRect(QRectF(0.5, 0.5, size - 1, size - 1), radius, radius)
        p.drawPath(path)
        pen = QPen(QColor("#ffffff"))
        pen.setWidth(3)
        pen.setCapStyle(Qt.PenCapStyle.RoundCap)
        pen.setJoinStyle(Qt.PenJoinStyle.RoundJoin)
        p.setPen(pen)
        p.setBrush(Qt.BrushStyle.NoBrush)
        p.drawLine(int(size * 0.22), int(size * 0.55), int(size * 0.44), int(size * 0.76))
        p.drawLine(int(size * 0.40), int(size * 0.76), int(size * 0.80), int(size * 0.26))
    else:
        p.setBrush(QBrush(QColor(COLORS['bg'])))
        pen = QPen(QColor(COLORS['border']))
        pen.setWidth(2)
        p.setPen(pen)
        path = QPainterPath()
        path.addRoundedRect(QRectF(1.0, 1.0, size - 2, size - 2), radius, radius)
        p.drawPath(path)
    p.end()
    return pm


def _ext_aware_path(path: str, new_ext: str) -> str:
    return os.path.splitext(path)[0] + new_ext


def _is_image_input(path: str) -> bool:
    ext = os.path.splitext(path)[1].lower().lstrip(".")
    return ext in EXTRA_INPUT_EXTS


# ============================================================
# 引擎探测
# ============================================================
def _find_pandoc() -> Optional[str]:
    try:
        return shutil.which("pandoc")
    except Exception:
        return None


def _find_wkhtmltopdf() -> Optional[str]:
    """探测 wkhtmltopdf 二进制（HTML→PDF 专用，体积约 50MB）"""
    for name in ("wkhtmltopdf", "wkhtmltopdf.exe"):
        try:
            p = shutil.which(name)
            if p:
                return p
        except Exception:
            pass
    for hard in (
        "/usr/bin/wkhtmltopdf",
        "/usr/local/bin/wkhtmltopdf",
        "/opt/homebrew/bin/wkhtmltopdf",
        r"C:\Program Files\wkhtmltopdf\bin\wkhtmltopdf.exe",
    ):
        if os.path.isfile(hard):
            return hard
    return None


def _probe_native_engines() -> dict:
    """探测当前 Python 环境内的文档处理库可用性（返回 flag 字典）"""
    flags = {}
    try:
        import docx2txt  # noqa: F401
        flags["docx2txt"] = True
    except Exception:
        flags["docx2txt"] = False
    try:
        from docx import Document  # noqa: F401
        flags["python_docx"] = True
    except Exception:
        flags["python_docx"] = False
    try:
        import openpyxl  # noqa: F401
        flags["openpyxl"] = True
    except Exception:
        flags["openpyxl"] = False
    try:
        import fitz  # PyMuPDF
        flags["fitz"] = True
    except Exception:
        flags["fitz"] = False
    try:
        import pypdf  # noqa: F401
        flags["pypdf"] = True
    except Exception:
        flags["pypdf"] = False
    try:
        from PIL import Image  # noqa: F401
        flags["PIL"] = True
    except Exception:
        flags["PIL"] = False
    try:
        import striprtf  # noqa: F401
        flags["striprtf"] = True
    except Exception:
        flags["striprtf"] = False
    try:
        import markdown  # noqa: F401
        flags["markdown"] = True
    except Exception:
        flags["markdown"] = False
    try:
        import csv as _csv  # noqa: F401
        flags["csv"] = True
    except Exception:
        flags["csv"] = False
    # ── 轻量文档扩展库 ──
    try:
        import pptx  # python-pptx
        flags["python_pptx"] = True
    except Exception:
        flags["python_pptx"] = False
    try:
        from pdf2docx import parse  # noqa: F401
        flags["pdf2docx"] = True
    except Exception:
        flags["pdf2docx"] = False
    try:
        from weasyprint import HTML  # noqa: F401
        flags["weasyprint"] = True
    except Exception:
        flags["weasyprint"] = False
    try:
        import pypandoc  # noqa: F401
        flags["pypandoc"] = True
    except Exception:
        flags["pypandoc"] = False
    return flags


# ============================================================
# Worker 线程
# ============================================================
class DocConversionWorker(QThread):
    progress_signal = Signal(int, int)
    task_started_signal = Signal(str)
    task_finished_signal = Signal(object)
    all_done_signal = Signal()
    single_progress_signal = Signal(int)

    def __init__(
        self,
        tasks: List[DocConversionTask],
        pandoc_path: Optional[str] = None,
        wkhtmltopdf_path: Optional[str] = None,
        native_flags: Optional[dict] = None,
    ):
        super().__init__()
        self.tasks = tasks
        self.pandoc_path = pandoc_path                # 主力引擎
        self.wkhtmltopdf_path = wkhtmltopdf_path      # HTML→PDF 专用
        self.native = native_flags or {}
        self._is_running = True

    # ----------------------------------------------------------
    # 1) 原生 Python 引擎（格式特定，速度最快）
    # ----------------------------------------------------------
    def _try_native_python(self, task: DocConversionTask) -> Tuple[bool, str]:
        """True=成功  False=回退；返回 (success, message)"""
        src, dst, fmt = task.input_path, task.output_path, task.output_format.lower()
        src_ext = os.path.splitext(src)[1].lower().lstrip(".")

        # --- 纯文本互转 (txt / md / html / csv / ini / py 等) ---
        pure_text_srcs = {"txt", "md", "markdown", "html", "htm", "rtf", "csv", "tsv",
                          "json", "xml", "yaml", "yml", "ini", "log", "py", "js", "ts",
                          "c", "cpp", "h", "mdx", "yaml"}
        pure_text_dsts = {"txt", "md", "html", "htm", "csv", "tsv"}
        if (src_ext in pure_text_srcs or fmt in pure_text_dsts) and not (src_ext.startswith("rtf") and fmt not in {"txt"}):
            # RTF → 文本，走 striprtf
            if src_ext == "rtf" and fmt == "txt" and self.native.get("striprtf"):
                try:
                    from striprtf.striprtf import rtf_to_text
                    with open(src, "r", encoding="utf-8", errors="replace") as f:
                        content = f.read()
                    text = rtf_to_text(content)
                    with open(dst, "w", encoding="utf-8") as fo:
                        fo.write(text)
                    return True, "RTF → TXT (striprtf)"
                except Exception as e:
                    return False, f"striprtf 失败: {e}"
            # Markdown → HTML
            if src_ext in {"md", "markdown"} and fmt == "html" and self.native.get("markdown"):
                try:
                    import markdown as md_lib
                    with open(src, "r", encoding="utf-8", errors="replace") as f:
                        html_body = md_lib.markdown(f.read(), extensions=["fenced_code", "tables"])
                    page = (
                        "<!DOCTYPE html><html><head><meta charset='utf-8'></head>"
                        f"<body style='font-family: sans-serif; line-height: 1.6;'>{html_body}</body></html>"
                    )
                    with open(dst, "w", encoding="utf-8") as fo:
                        fo.write(page)
                    return True, "Markdown → HTML (python-markdown)"
                except Exception as e:
                    return False, f"python-markdown 失败: {e}"
            # CSV → XLSX （若 openpyxl 存在）
            if src_ext == "csv" and fmt == "xlsx" and self.native.get("openpyxl"):
                try:
                    import csv
                    from openpyxl import Workbook
                    wb = Workbook()
                    ws = wb.active
                    ws.title = "Sheet1"
                    with open(src, "r", encoding="utf-8-sig", errors="replace", newline="") as f:
                        reader = csv.reader(f)
                        for i, row in enumerate(reader, start=1):
                            for j, cell in enumerate(row, start=1):
                                ws.cell(row=i, column=j, value=cell)
                    save_dir = os.path.dirname(dst)
                    if save_dir:
                        os.makedirs(save_dir, exist_ok=True)
                    wb.save(dst)
                    return True, "CSV → XLSX (openpyxl)"
                except Exception as e:
                    return False, f"openpyxl 失败: {e}"
            # XLSX → CSV （openpyxl）
            if src_ext == "xlsx" and fmt == "csv" and self.native.get("openpyxl"):
                try:
                    import csv
                    from openpyxl import load_workbook
                    wb = load_workbook(src, read_only=True, data_only=True)
                    ws = wb[wb.sheetnames[0]]
                    save_dir = os.path.dirname(dst)
                    if save_dir:
                        os.makedirs(save_dir, exist_ok=True)
                    with open(dst, "w", encoding="utf-8-sig", newline="") as fo:
                        writer = csv.writer(fo)
                        for row in ws.iter_rows(values_only=True):
                            writer.writerow(["" if c is None else str(c) for c in row])
                    return True, "XLSX → CSV (openpyxl)"
                except Exception as e:
                    return False, f"openpyxl 失败: {e}"
            # DOCX → TXT （docx2txt 优先，否则 python-docx 兜底）
            if src_ext == "docx" and fmt == "txt":
                try:
                    if self.native.get("docx2txt"):
                        import docx2txt
                        text = docx2txt.process(src) or ""
                        with open(dst, "w", encoding="utf-8") as fo:
                            fo.write(text)
                        return True, "DOCX → TXT (docx2txt)"
                    if self.native.get("python_docx"):
                        from docx import Document
                        doc = Document(src)
                        paragraphs = [p.text for p in doc.paragraphs]
                        with open(dst, "w", encoding="utf-8") as fo:
                            fo.write("\n".join(paragraphs))
                        return True, "DOCX → TXT (python-docx)"
                except Exception as e:
                    return False, f"docx 读取失败: {e}"
            # 普通纯文本 → 纯文本（复制/编码转换）：目标或源是纯文本时
            if (src_ext in pure_text_srcs) and (fmt in pure_text_dsts):
                try:
                    with open(src, "rb") as f:
                        raw = f.read()
                    text = raw.decode("utf-8", errors="replace")
                    save_dir = os.path.dirname(dst)
                    if save_dir:
                        os.makedirs(save_dir, exist_ok=True)
                    with open(dst, "w", encoding="utf-8") as fo:
                        fo.write(text)
                    return True, "文本复制 + UTF-8 标准化"
                except Exception as e:
                    return False, f"文本读写失败: {e}"

        # --- 图片 → PDF（Pillow） ---
        if _is_image_input(src) and fmt == "pdf" and self.native.get("PIL"):
            try:
                from PIL import Image
                with Image.open(src) as img:
                    img.load()
                    if img.mode in ("RGBA", "P", "LA"):
                        img = img.convert("RGB")
                    save_dir = os.path.dirname(dst)
                    if save_dir:
                        os.makedirs(save_dir, exist_ok=True)
                    img.save(dst, "PDF", resolution=task.pdf_dpi)
                    return True, f"图片 → PDF (Pillow {img.size[0]}x{img.size[1]} DPI={task.pdf_dpi})"
            except Exception as e:
                return False, f"Pillow 图片→PDF 失败: {e}"

        # --- PDF → 图片（PyMuPDF） ---
        if src_ext == "pdf" and _is_image_input(dst):  # dst 按扩展名当图片单页导出
            if self.native.get("fitz"):
                try:
                    import fitz
                    doc = fitz.open(src)
                    save_dir = os.path.dirname(dst)
                    if save_dir:
                        os.makedirs(save_dir, exist_ok=True)
                    # 导出第 1 页
                    if doc.page_count >= 1:
                        page = doc.load_page(0)
                        mat = fitz.Matrix(task.pdf_dpi / 72.0, task.pdf_dpi / 72.0)
                        pix = page.get_pixmap(matrix=mat, alpha=False)
                        base, ext = os.path.splitext(dst)
                        if doc.page_count == 1:
                            pix.save(dst)
                        else:
                            # 多页：加 _p1 / _p2 编号
                            for i in range(doc.page_count):
                                p = doc.load_page(i)
                                pix = p.get_pixmap(matrix=mat, alpha=False)
                                out = f"{base}_p{i+1}{ext}"
                                pix.save(out)
                    doc.close()
                    return True, f"PDF → 图片 (PyMuPDF DPI={task.pdf_dpi}, pages={doc.page_count})"
                except Exception as e:
                    return False, f"PyMuPDF 失败: {e}"

        # --- PDF 文本提取 (pypdf / fitz) ---
        if src_ext == "pdf" and fmt == "txt":
            if self.native.get("fitz"):
                try:
                    import fitz
                    doc = fitz.open(src)
                    texts = []
                    for i in range(doc.page_count):
                        texts.append(doc.load_page(i).get_text("text") or "")
                    doc.close()
                    with open(dst, "w", encoding="utf-8") as fo:
                        fo.write("\n\n".join(texts))
                    return True, "PDF → TXT (PyMuPDF)"
                except Exception:
                    pass
            if self.native.get("pypdf"):
                try:
                    from pypdf import PdfReader
                    reader = PdfReader(src)
                    with open(dst, "w", encoding="utf-8") as fo:
                        for page in reader.pages:
                            try:
                                fo.write(page.extract_text() or "")
                            except Exception:
                                pass
                            fo.write("\n\n")
                    return True, "PDF → TXT (pypdf)"
                except Exception as e:
                    return False, f"pypdf 失败: {e}"
        return False, "当前格式组合无原生引擎，尝试外部工具…"

    # ----------------------------------------------------------
    # 1.5) pdf2docx — PDF → DOCX（替代 LO 的 PDF→DOCX）
    # ----------------------------------------------------------
    def _try_pdf2docx(self, task: DocConversionTask) -> Tuple[bool, str]:
        if not self.native.get("pdf2docx"):
            return False, "未安装 pdf2docx"
        src, dst, fmt = task.input_path, task.output_path, task.output_format.lower()
        src_ext = os.path.splitext(src)[1].lower().lstrip(".")
        if src_ext != "pdf" or fmt != "docx":
            return False, "仅支持 PDF → DOCX"
        save_dir = os.path.dirname(dst)
        if save_dir:
            os.makedirs(save_dir, exist_ok=True)
        try:
            from pdf2docx import parse
            parse(src, dst, start=0, end=None, multi_processing=False)
            if os.path.isfile(dst) and os.path.getsize(dst) > 0:
                return True, "PDF → DOCX (pdf2docx)"
            return False, "pdf2docx 输出文件为空"
        except Exception as e:
            return False, f"pdf2docx 失败: {e}"

    # ----------------------------------------------------------
    # 1.6) python-pptx — PPTX 提取文本
    # ----------------------------------------------------------
    def _try_python_pptx(self, task: DocConversionTask) -> Tuple[bool, str]:
        if not self.native.get("python_pptx"):
            return False, "未安装 python-pptx"
        src, dst, fmt = task.input_path, task.output_path, task.output_format.lower()
        src_ext = os.path.splitext(src)[1].lower().lstrip(".")
        if src_ext not in ("pptx", "ppt") or fmt not in ("txt", "pdf"):
            return False, "仅支持 PPTX → TXT / PDF"
        save_dir = os.path.dirname(dst)
        if save_dir:
            os.makedirs(save_dir, exist_ok=True)
        try:
            from pptx import Presentation

            if src_ext == "ppt":
                # python-pptx 不支持老 .ppt 二进制格式
                return False, "python-pptx 不支持 .ppt 老格式"

            prs = Presentation(src)

            # ── PPTX → TXT：提取所有幻灯片文本 ──
            if fmt == "txt":
                texts = []
                for i, slide in enumerate(prs.slides, start=1):
                    texts.append(f"===== Slide {i} =====")
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            for para in shape.text_frame.paragraphs:
                                line = "".join(run.text for run in para.runs)
                                if line:
                                    texts.append(line)
                with open(dst, "w", encoding="utf-8") as fo:
                    fo.write("\n".join(texts))
                return True, f"PPTX → TXT (python-pptx, {len(prs.slides)} 页)"

            # ── PPTX → PDF：python-pptx 无原生渲染能力 ──
            if fmt == "pdf":
                return False, "python-pptx 无原生 PDF 渲染能力，暂不支持 PPTX→PDF"
        except Exception as e:
            return False, f"python-pptx 失败: {e}"

    # ----------------------------------------------------------
    # 1.7) HTML → PDF（weasyprint 或 wkhtmltopdf 二选一）
    # ----------------------------------------------------------
    def _try_html_to_pdf(self, task: DocConversionTask) -> Tuple[bool, str]:
        src, dst, fmt = task.input_path, task.output_path, task.output_format.lower()
        src_ext = os.path.splitext(src)[1].lower().lstrip(".")
        if src_ext not in ("html", "htm") or fmt != "pdf":
            return False, "仅支持 HTML → PDF"
        save_dir = os.path.dirname(dst)
        if save_dir:
            os.makedirs(save_dir, exist_ok=True)

        # 优先 weasyprint（纯 Python，无外部依赖）
        if self.native.get("weasyprint"):
            try:
                from weasyprint import HTML
                HTML(filename=src).write_pdf(dst)
                if os.path.isfile(dst) and os.path.getsize(dst) > 0:
                    return True, "HTML → PDF (weasyprint)"
            except Exception:
                # weasyprint 失败时降级到 wkhtmltopdf
                pass

        # 次选 wkhtmltopdf（外部二进制，需系统安装）
        if self.wkhtmltopdf_path:
            try:
                cmd = [
                    self.wkhtmltopdf_path,
                    "--quiet",
                    "--encoding", "utf-8",
                    "--enable-local-file-access",
                    src, dst,
                ]
                result = subprocess.run(
                    cmd, capture_output=True, text=True, timeout=120,
                    encoding="utf-8", errors="replace"
                )
                if result.returncode == 0 and os.path.isfile(dst) and os.path.getsize(dst) > 0:
                    return True, "HTML → PDF (wkhtmltopdf)"
                tail = "\n".join(deque(result.stderr.splitlines(), maxlen=6)) or "(无输出)"
                return False, f"wkhtmltopdf 失败 (exit={result.returncode})\n{tail}"
            except Exception as e:
                return False, f"wkhtmltopdf 异常: {e}"

        return False, "未检测到 weasyprint 或 wkhtmltopdf"

    # ----------------------------------------------------------
    # 2) Pandoc 主力引擎（通用文档互转）
    # ----------------------------------------------------------
    def _try_pandoc(self, task: DocConversionTask) -> Tuple[bool, str]:
        """pandoc 主引擎：处理 MD / HTML / DOCX / ODT / EPUB / RTF / LaTeX 等互转"""
        if not self.pandoc_path:
            return False, "未检测到 pandoc"
        src, dst, fmt = task.input_path, task.output_path, task.output_format.lower()
        src_ext = os.path.splitext(src)[1].lower().lstrip(".")
        save_dir = os.path.dirname(dst)
        if save_dir:
            os.makedirs(save_dir, exist_ok=True)
        # pandoc 支持的输入格式
        pandoc_inputs = {
            "md", "markdown", "html", "htm", "docx", "odt", "epub",
            "rst", "latex", "tex", "rtf", "txt", "json", "mediawiki",
            "t2t", "org", "wiki",
        }
        # pandoc 支持的输出格式
        pandoc_outputs = {
            "md", "markdown", "html", "htm", "docx", "odt", "epub",
            "rst", "latex", "tex", "rtf", "txt", "json", "mediawiki",
            "pdf",  # 需配合 LaTeX 引擎（如 tectonic / xelatex）
        }
        if src_ext not in pandoc_inputs or fmt not in pandoc_outputs:
            return False, "格式不在 pandoc 支持范围"
        cmd = [self.pandoc_path, src, "-o", dst, "--standalone"]
        # ── PDF 输出：pandoc 需要配合 LaTeX 引擎 ──
        if fmt == "pdf":
            tectonic = shutil.which("tectonic")
            if tectonic:
                cmd.extend(["--pdf-engine=tectonic"])
            else:
                xelatex = shutil.which("xelatex")
                if xelatex:
                    cmd.extend(["--pdf-engine=xelatex"])
                else:
                    return False, "pandoc → PDF 需安装 tectonic 或 xelatex"
        # ── 中文字体处理（DOCX/HTML/EPUB 输出时） ──
        if fmt in ("docx", "html", "epub"):
            cmd.append("--metadata=lang:zh-CN")
        try:
            result = subprocess.run(
                cmd, capture_output=True, text=True, timeout=120,
                encoding="utf-8", errors="replace"
            )
            if result.returncode == 0 and os.path.isfile(dst) and os.path.getsize(dst) > 0:
                return True, f"pandoc → {fmt.upper()}"
            tail = "\n".join(deque(result.stderr.splitlines(), maxlen=8)) or "(无输出)"
            return False, f"pandoc 失败 (exit={result.returncode})\n--- stderr ---\n{tail}"
        except Exception as e:
            return False, f"pandoc 异常: {e}"

    # ----------------------------------------------------------
    # 主入口
    # ----------------------------------------------------------
    def run(self):
        total = len(self.tasks)
        engine_hint = []
        if self.native:
            enabled = [k for k, v in self.native.items() if v]
            if enabled:
                engine_hint.append("原生Python库: " + ",".join(enabled))
        if self.pandoc_path:
            engine_hint.append(f"pandoc: {os.path.basename(self.pandoc_path)}")
        if self.wkhtmltopdf_path:
            engine_hint.append(f"wkhtmltopdf: {os.path.basename(self.wkhtmltopdf_path)}")
        ok = 0
        fail = 0
        for idx, task in enumerate(self.tasks, start=1):
            if not self._is_running:
                break
            name = os.path.basename(task.input_path)
            self.task_started_signal.emit(name)
            self.single_progress_signal.emit(5)
            QThread.msleep(10)
            self.single_progress_signal.emit(30)

            save_dir = os.path.dirname(task.output_path)
            if save_dir and not os.path.isdir(save_dir):
                os.makedirs(save_dir, exist_ok=True)

            used = []
            res_msg = ""
            success = False
            last_msg = ""

            def _push(label, s, m):
                nonlocal success, res_msg, last_msg
                used.append(label)
                last_msg = m
                if s:
                    success = True
                    res_msg = m

            # 1) 原生 Python 引擎（含原有库 + 新增库）
            s, m = self._try_native_python(task)
            _push("原生Python", s, m)
            self.single_progress_signal.emit(40)

            # 2) PDF → DOCX 专用：pdf2docx
            if not success and self._is_running:
                s, m = self._try_pdf2docx(task)
                _push("pdf2docx", s, m)
            self.single_progress_signal.emit(55)

            # 3) PPTX 专用：python-pptx
            if not success and self._is_running:
                s, m = self._try_python_pptx(task)
                _push("python-pptx", s, m)
            self.single_progress_signal.emit(65)

            # 4) pandoc 主力（替换 LO 的通用互转）
            if not success and self._is_running:
                s, m = self._try_pandoc(task)
                _push("pandoc", s, m)
            self.single_progress_signal.emit(80)

            # 5) HTML → PDF 专用（weasyprint / wkhtmltopdf）
            if not success and self._is_running:
                s, m = self._try_html_to_pdf(task)
                _push("HTML→PDF", s, m)
            self.single_progress_signal.emit(100)

            result = DocConversionResult(success, task, res_msg if success else (last_msg or "全部引擎失败"))
            self.task_finished_signal.emit(result)
            if success:
                ok += 1
            else:
                fail += 1
            self.progress_signal.emit(idx, total)

        self.all_done_signal.emit()

    def cancel(self):
        self._is_running = False


# ============================================================
# UI 主 Widget
# ============================================================
class DocConverterWidget(QWidget):
    task_monitor_signal = Signal(str)
    log_signal = Signal(str, str)

    def __init__(self, default_output_dir: str = ""):
        super().__init__()
        self.worker = None
        self.theme_colors = ThemeManager.instance().current_colors
        self._default_output_dir = (default_output_dir or "").strip()

        # 引擎探测
        self.native_flags: dict = _probe_native_engines()
        self.pandoc_path: Optional[str] = _find_pandoc()
        self.wkhtmltopdf_path: Optional[str] = _find_wkhtmltopdf()

        self._setup_ui()
        self._apply_widget_styles()
        if self._default_output_dir and not self.output_path_edit.text().strip():
            self.output_path_edit.setText(self._default_output_dir)
        try:
            ThemeManager.instance().theme_changed.connect(self.reapply_theme)
        except Exception:
            pass
        QTimer.singleShot(200, self._announce_engines)

    # ============================================================
    # UI 构建
    # ============================================================
    def _setup_ui(self):
        self.setMinimumWidth(700)
        self.resize(1000, 700)

        layout = QVBoxLayout(self)
        layout.setContentsMargins(20, 8, 20, 20)
        layout.setSpacing(14)

        title_label = QLabel("📄  文档格式转换")
        title_label.setStyleSheet(f"font-size: 22px; font-weight: bold; color: {COLORS['text']};")
        subtitle_label = QLabel("支持 PDF / Word / Excel / PPT / Markdown / HTML / EPUB 等批量互转（Python原生库 + pandoc + pdf2docx + python-pptx 轻量引擎链）")
        subtitle_label.setStyleSheet(f"font-size: 13px; color: {COLORS['text_secondary']};")
        title_layout = QVBoxLayout()
        title_layout.addWidget(title_label)
        title_layout.addWidget(subtitle_label)
        title_layout.setSpacing(4)
        layout.addLayout(title_layout)

        top_widget = QWidget()
        layout.addWidget(top_widget, 1)
        top_layout = QHBoxLayout(top_widget)
        top_layout.setContentsMargins(0, 0, 0, 0)
        top_layout.setSpacing(16)

        # ------- 左：待转换文件 -------
        self._build_file_panel(top_layout)
        # ------- 右：转换设置 -------
        self._build_settings_panel(top_layout)

        # 下半：进度 + 日志 + 开始/停止
        bottom_widget = QWidget()
        self._build_progress_panel(bottom_widget)
        layout.addWidget(bottom_widget)

    def _build_file_panel(self, top_layout: QHBoxLayout):
        file_group = QGroupBox("待转换文档")
        file_group.setStyleSheet(f"""
            QGroupBox {{
                font-weight: bold;
                color: {COLORS['text']};
                border: 1px solid {COLORS['border']};
                border-radius: 8px;
                margin-top: 10px;
                padding-top: 16px;
                background-color: {COLORS['card']};
            }}
            QGroupBox::title {{
                subcontrol-origin: margin;
                left: 12px;
                padding: 0 6px;
            }}
        """)
        file_layout = QVBoxLayout(file_group)
        file_layout.setContentsMargins(12, 16, 12, 12)

        self.file_list = QListWidget()
        self.file_list.setSelectionMode(QListWidget.SelectionMode.ExtendedSelection)
        self.file_list.setAcceptDrops(True)
        self.file_list.dragEnterEvent = self._drag_enter
        self.file_list.dropEvent = self._drop
        self.file_list.setContextMenuPolicy(Qt.ContextMenuPolicy.CustomContextMenu)
        self.file_list.customContextMenuRequested.connect(self._show_file_context_menu)
        self.file_list.setMinimumHeight(280)
        self.file_list.setStyleSheet(f"""
            QListWidget {{
                background-color: {COLORS['bg']};
                border: 1px solid {COLORS['border']};
                border-radius: 6px;
                color: {COLORS['text']};
                padding: 4px;
            }}
            QListWidget::item {{
                padding: 8px;
                border-bottom: 1px solid {COLORS['border']};
            }}
            QListWidget::item:selected {{
                background-color: {COLORS['primary']};
                color: white;
            }}
        """)
        file_layout.addWidget(self.file_list, 1)

        file_btn_layout = QHBoxLayout()
        self.add_btn = QPushButton("➕ 添加文件")
        self.add_btn.clicked.connect(self._add_files)
        self.remove_btn = QPushButton("➖ 移除选中")
        self.remove_btn.clicked.connect(self._remove_selected)
        self.clear_btn = QPushButton("🗑 清空全部")
        self.clear_btn.clicked.connect(self._clear_files)
        for btn in [self.add_btn, self.remove_btn, self.clear_btn]:
            btn.setStyleSheet(f"""
                QPushButton {{
                    background-color: {COLORS['card_hover']};
                    color: {COLORS['text']};
                    border: 1px solid {COLORS['border']};
                    border-radius: 4px;
                    padding: 6px 12px;
                    font-size: 13px;
                }}
                QPushButton:hover {{
                    background-color: {COLORS['primary']};
                    color: white;
                    border-color: {COLORS['primary']};
                }}
            """)
            file_btn_layout.addWidget(btn)
        file_layout.addLayout(file_btn_layout)
        top_layout.addWidget(file_group, 3)

    def _build_settings_panel(self, top_layout: QHBoxLayout):
        settings_group = QGroupBox("转换设置")
        settings_group.setMinimumWidth(360)
        settings_group.setStyleSheet(f"""
            QGroupBox {{
                font-weight: bold;
                color: {COLORS['text']};
                border: 1px solid {COLORS['border']};
                border-radius: 8px;
                margin-top: 10px;
                padding-top: 4px;
                background-color: {COLORS['card']};
            }}
            QGroupBox::title {{
                subcontrol-origin: margin;
                left: 12px;
                padding: 0 6px;
            }}
        """)
        settings_layout = QVBoxLayout(settings_group)
        settings_layout.setContentsMargins(18, 6, 18, 14)
        settings_layout.setSpacing(10)
        settings_layout.setSizeConstraint(settings_layout.SizeConstraint.SetMinimumSize)

        input_style = f"""
            QComboBox, QLineEdit, QSpinBox {{
                background-color: {COLORS['bg']};
                color: {COLORS['text']};
                border: 1px solid {COLORS['border']};
                border-radius: 4px;
                padding: 6px 10px;
                font-size: 13px;
                min-height: 22px;
            }}
            QComboBox:hover, QLineEdit:hover, QSpinBox:hover {{
                border-color: {COLORS['primary']};
            }}
            QComboBox:focus, QLineEdit:focus, QSpinBox:focus {{
                border-color: {COLORS['success']};
            }}
            QComboBox::drop-down {{
                border: none;
                width: 20px;
            }}
            QSpinBox::up-button, QSpinBox::down-button {{
                width: 16px;
            }}
        """

        def create_setting_row(label_text, right_widget):
            row_widget = QWidget()
            row_widget.setMinimumHeight(44)
            row_widget.setSizePolicy(
                QSizePolicy.Policy.Preferred, QSizePolicy.Policy.Fixed
            )
            row = QHBoxLayout(row_widget)
            row.setContentsMargins(0, 0, 0, 0)
            row.setSpacing(12)

            lab = QLabel(label_text)
            lab.setStyleSheet(
                get_cjk_font_qss(
                    font_size_px=14,
                    color=COLORS['text'],
                    extra="font-weight: 400; background: transparent;"
                )
            )
            lab.setMinimumWidth(80)
            lab.setSizePolicy(
                QSizePolicy.Policy.MinimumExpanding, QSizePolicy.Policy.Preferred
            )
            lab.setAlignment(Qt.AlignmentFlag.AlignRight | Qt.AlignmentFlag.AlignVCenter)

            row.addWidget(lab, 0)
            row.addWidget(right_widget, 1)
            row.setStretch(0, 1)
            row.setStretch(1, 3)
            return row_widget

        # 输出格式
        self.format_combo = QComboBox()
        for key, info in SUPPORTED_DOC_FORMATS.items():
            self.format_combo.addItem(info["desc"], key)
        # 默认 PDF
        self.format_combo.setCurrentIndex(list(SUPPORTED_DOC_FORMATS.keys()).index("pdf"))
        self.format_combo.setStyleSheet(input_style)
        settings_layout.addWidget(create_setting_row("输出格式", self.format_combo))

        # 输出目录
        self.output_path_edit = QLineEdit()
        self.output_path_edit.setPlaceholderText("留空 = 在源文件同目录输出")
        self.output_path_edit.setStyleSheet(input_style)
        self.output_browse_btn = QPushButton("📂 浏览")
        self.output_browse_btn.setStyleSheet(f"""
            QPushButton {{
                background-color: {COLORS['card_hover']};
                color: {COLORS['text']};
                border: 1px solid {COLORS['border']};
                border-radius: 4px;
                padding: 6px 10px;
            }}
            QPushButton:hover {{
                background-color: {COLORS['primary']};
                color: white;
                border-color: {COLORS['primary']};
            }}
        """)
        self.output_browse_btn.clicked.connect(self._browse_output_dir)
        out_wrapper = QWidget()
        out_wrapper.setMinimumHeight(36)
        out_wrap_layout = QHBoxLayout(out_wrapper)
        out_wrap_layout.setContentsMargins(0, 0, 0, 0)
        out_wrap_layout.setSpacing(8)
        out_wrap_layout.addWidget(self.output_path_edit, 1)
        out_wrap_layout.addWidget(self.output_browse_btn, 0)
        settings_layout.addWidget(create_setting_row("输出目录", out_wrapper))

        # PDF 导出 DPI
        self.dpi_spin = QSpinBox()
        self.dpi_spin.setRange(72, 600)
        self.dpi_spin.setSingleStep(12)
        self.dpi_spin.setValue(200)
        self.dpi_spin.setSuffix(" DPI")
        self.dpi_spin.setStyleSheet(input_style)
        settings_layout.addWidget(create_setting_row("PDF 渲染 DPI", self.dpi_spin))

        # 保留原文件结构（文件夹输入时） 复用 overwrite 模式
        self.overwrite_check = QPushButton(
            QIcon(_make_checkbox_icon_pixmap(False, COLORS['accent'])),
            "  重名文件 — 直接覆盖 (否则自动重命名)"
        )
        self.overwrite_check.setCheckable(True)
        self.overwrite_check.setChecked(False)
        self.overwrite_check.setIconSize(QSize(22, 22))
        self.overwrite_check.setCursor(Qt.PointingHandCursor)
        self.overwrite_check.setMinimumHeight(38)
        self.overwrite_check.setStyleSheet(f"""
            QPushButton {{
                color: {COLORS['text']};
                font-size: 14px;
                font-weight: 500;
                padding: 6px 12px 6px 10px;
                border-radius: 6px;
                border: none;
                background-color: transparent;
                text-align: left;
            }}
            QPushButton:hover {{
                background-color: rgba(233, 69, 96, 0.10);
            }}
            QPushButton:checked {{
                background-color: rgba(233, 69, 96, 0.18);
                color: #ffffff;
                font-weight: 600;
            }}
        """)
        def _on_over_toggled(c):
            self.overwrite_check.setIcon(QIcon(_make_checkbox_icon_pixmap(c, COLORS['accent'])))
        self.overwrite_check.toggled.connect(_on_over_toggled)
        cb2 = QHBoxLayout()
        cb2.addSpacing(108)
        cb2.addWidget(self.overwrite_check)
        cb2.addStretch()
        settings_layout.addLayout(cb2)

        # 引擎状态面板（只读 label，启动后显示探测结果）
        self.engine_status_label = QLabel("引擎状态：加载中…")
        self.engine_status_label.setStyleSheet(
            f"color: {COLORS['text_secondary']};"
            "font-size: 12px;"
            "padding: 6px 10px;"
            f"background-color: {COLORS['bg']};"
            f"border: 1px solid {COLORS['border']};"
            "border-radius: 6px;"
        )
        self.engine_status_label.setWordWrap(True)
        settings_layout.addWidget(self.engine_status_label)

        settings_layout.addStretch(1)
        top_layout.addWidget(settings_group, 2)

    def _build_progress_panel(self, bottom_widget: QWidget):
        outer = QVBoxLayout(bottom_widget)
        outer.setContentsMargins(0, 0, 0, 0)
        outer.setSpacing(10)

        # 开始/停止按钮
        btn_row = QHBoxLayout()
        btn_row.setSpacing(12)
        self.convert_btn = QPushButton("▶  开始转换")
        self.convert_btn.clicked.connect(self._start_conversion)
        self.convert_btn.setMinimumHeight(44)
        self.convert_btn.setStyleSheet(f"""
            QPushButton {{
                background-color: {COLORS['success']};
                color: white;
                border: none;
                border-radius: 8px;
                font-size: 15px;
                font-weight: 600;
            }}
            QPushButton:hover {{
                background-color: {COLORS['primary']};
            }}
            QPushButton:disabled {{
                background-color: {COLORS['card_hover']};
                color: {COLORS['text_secondary']};
            }}
        """)
        self.stop_btn = QPushButton("⏹  停止")
        self.stop_btn.clicked.connect(self._stop_conversion)
        self.stop_btn.setEnabled(False)
        self.stop_btn.setMinimumHeight(44)
        self.stop_btn.setStyleSheet(f"""
            QPushButton {{
                background-color: {COLORS['error']};
                color: white;
                border: none;
                border-radius: 8px;
                font-size: 15px;
                font-weight: 600;
            }}
            QPushButton:hover {{ background-color: #ff6b6b; }}
            QPushButton:disabled {{
                background-color: {COLORS['card_hover']};
                color: {COLORS['text_secondary']};
            }}
        """)
        self.open_dir_btn = QPushButton("📂  打开输出目录")
        self.open_dir_btn.clicked.connect(self._open_output_dir)
        self.open_dir_btn.setMinimumHeight(44)
        self.open_dir_btn.setStyleSheet(f"""
            QPushButton {{
                background-color: {COLORS['card_hover']};
                color: {COLORS['text']};
                border: 1px solid {COLORS['border']};
                border-radius: 8px;
                font-size: 15px;
                font-weight: 600;
            }}
            QPushButton:hover {{
                background-color: {COLORS['primary']};
                color: white;
                border-color: {COLORS['primary']};
            }}
        """)
        btn_row.addWidget(self.convert_btn, 3)
        btn_row.addWidget(self.stop_btn, 1)
        btn_row.addWidget(self.open_dir_btn, 2)
        outer.addLayout(btn_row)

    # ============================================================
    # 引擎状态
    # ============================================================
    def _announce_engines(self):
        lines = []
        native_hits = [k for k, v in self.native_flags.items() if v]
        if native_hits:
            lines.append("✅ Python 原生库: " + ",".join(native_hits))
        else:
            lines.append("❌ 无可用 Python 原生文档库（缺少依赖，请重新安装完整包）")

        if self.pandoc_path:
            lines.append(f"✅ pandoc (主力): {self.pandoc_path}")
        else:
            lines.append("⚠️ 未检测到 pandoc（缺少依赖，请重新安装完整包）")

        if self.wkhtmltopdf_path:
            lines.append(f"✅ wkhtmltopdf: {self.wkhtmltopdf_path}")
        elif not self.native_flags.get("weasyprint"):
            lines.append("ℹ️ 未检测到 HTML→PDF 引擎（推荐 weasyprint 或 wkhtmltopdf）")

        # 新增轻量扩展库状态
        new_libs = []
        if self.native_flags.get("python_pptx"):
            new_libs.append("python-pptx")
        if self.native_flags.get("pdf2docx"):
            new_libs.append("pdf2docx")
        if self.native_flags.get("weasyprint"):
            new_libs.append("weasyprint")
        if new_libs:
            lines.append("✅ 轻量扩展库: " + ",".join(new_libs))
        else:
            lines.append("ℹ️ 可选轻量库未随包附带（缺少依赖，请重新安装完整包）")

        status_html = "<br>".join(lines)
        self.engine_status_label.setText("引擎状态：<br>" + status_html)

    # ============================================================
    # 文件列表 & 拖拽
    # ============================================================
    def _drag_enter(self, event: QDragEnterEvent):
        if event.mimeData().hasUrls():
            event.acceptProposedAction()

    def _drop(self, event: QDropEvent):
        for url in event.mimeData().urls():
            path = url.toLocalFile()
            if not path:
                continue
            if os.path.isdir(path):
                for root, dirs, files in os.walk(path):
                    for fn in files:
                        ext = os.path.splitext(fn)[1].lower().lstrip(".")
                        if ext in ALL_INPUT_EXTS:
                            self._add_file(os.path.join(root, fn))
            else:
                self._add_file(path)
        event.acceptProposedAction()

    def _add_file(self, path):
        if not path or not os.path.isfile(path):
            return
        ext = os.path.splitext(path)[1].lower().lstrip(".")
        if ext not in ALL_INPUT_EXTS:
            # 仍允许加入，但提示
            pass
        # 去重
        for i in range(self.file_list.count()):
            if self.file_list.item(i).data(Qt.ItemDataRole.UserRole) == path:
                return
        item = QListWidgetItem(f"{os.path.basename(path)}   ({self._fmt_size(path)})")
        item.setData(Qt.ItemDataRole.UserRole, path)
        if ext in EXTRA_INPUT_EXTS:
            item.setToolTip("图片文件：支持转 PDF 或导出为图片格式（通过 DPI 选项）")
        self.file_list.addItem(item)

    def _fmt_size(self, path):
        try:
            s = os.path.getsize(path)
        except Exception:
            return "?"
        if s < 1024:
            return f"{s} B"
        if s < 1024 * 1024:
            return f"{s/1024:.1f} KB"
        return f"{s/(1024*1024):.2f} MB"

    def _add_files(self):
        all_ext_list = (
            " ".join(f"*.{e}" for e in sorted(ALL_INPUT_EXTS))
            + " *.DOCX *.XLSX *.PPTX *.PDF *.TXT *.MD *.HTML *.EPUB"
        )
        files, _ = QFileDialog.getOpenFileNames(
            self, "选择要转换的文档/图片",
            QStandardPaths.writableLocation(QStandardPaths.StandardLocation.DesktopLocation) or os.path.expanduser("~"),
            f"所有支持的格式 ({all_ext_list});;所有文件 (*)"
        )
        for f in files:
            self._add_file(f)

    def _remove_selected(self):
        for item in self.file_list.selectedItems():
            self.file_list.takeItem(self.file_list.row(item))

    def _clear_files(self):
        self.file_list.clear()

    def _show_file_context_menu(self, pos):
        menu = QMenu(self)
        act_add = menu.addAction("➕ 添加文件")
        act_rm = menu.addAction("➖ 移除选中")
        act_clear = menu.addAction("🗑 清空全部")
        chosen = menu.exec(self.file_list.mapToGlobal(pos))
        if chosen == act_add:
            self._add_files()
        elif chosen == act_rm:
            self._remove_selected()
        elif chosen == act_clear:
            self._clear_files()

    # ============================================================
    # 输出目录
    # ============================================================
    def _browse_output_dir(self):
        dir_ = QFileDialog.getExistingDirectory(
            self, "选择输出目录",
            self.output_path_edit.text()
            or QStandardPaths.writableLocation(QStandardPaths.StandardLocation.DesktopLocation)
            or os.path.expanduser("~")
        )
        if dir_:
            self.output_path_edit.setText(dir_)

    def _open_output_dir(self):
        path = None
        if self.output_path_edit.text().strip():
            path = self.output_path_edit.text().strip()
        else:
            # 找第一个文件的父目录
            if self.file_list.count() > 0:
                p = self.file_list.item(0).data(Qt.ItemDataRole.UserRole)
                if p:
                    path = os.path.dirname(p)
        if not path or not os.path.isdir(path):
            QMessageBox.information(self, "提示", "未指定输出目录且无文件，请先添加文件或设置输出目录。")
            return
        try:
            if sys.platform == 'win32':
                os.startfile(path)  # type: ignore[attr-defined]
            elif sys.platform == 'darwin':
                subprocess.run(['open', path])
            else:
                subprocess.run(['xdg-open', path])
        except Exception as e:
            QMessageBox.warning(self, "失败", f"无法打开目录: {e}")

    # ============================================================
    # 构建任务 / 启动转换
    # ============================================================
    def _build_tasks(self):
        tasks = []
        output_format = self.format_combo.currentData()
        output_dir = self.output_path_edit.text().strip()
        overwrite = self.overwrite_check.isChecked()
        pdf_dpi = int(self.dpi_spin.value())
        ext = SUPPORTED_DOC_FORMATS[output_format]["ext"]

        selected_items = self.file_list.selectedItems()
        items_to_convert = selected_items if selected_items else [
            self.file_list.item(i) for i in range(self.file_list.count())
        ]

        for item in items_to_convert:
            input_path = item.data(Qt.ItemDataRole.UserRole)
            if output_dir:
                base = os.path.splitext(os.path.basename(input_path))[0]
                out_path = os.path.join(output_dir, base + ext)
            else:
                base = os.path.splitext(input_path)[0]
                out_path = base + ext

            if not overwrite:
                counter = 1
                original = out_path
                while os.path.isfile(out_path):
                    out_path = os.path.splitext(original)[0] + f"_{counter}" + ext
                    counter += 1

            tasks.append(DocConversionTask(
                input_path=input_path,
                output_path=out_path,
                output_format=output_format,
                pdf_dpi=pdf_dpi,
                overwrite=overwrite,
            ))
        return tasks

    def _ensure_any_engine(self) -> bool:
        any_native = any(self.native_flags.values())
        if any_native or self.pandoc_path:
            return True
        dlg = QMessageBox(self)
        dlg.setIcon(QMessageBox.Icon.Warning)
        dlg.setWindowTitle("文档转换缺少依赖")
        dlg.setTextFormat(Qt.TextFormat.RichText)
        dlg.setText(
            "当前环境缺少文档转换引擎：<br><br>"
            "① Python 原生文档库（python-docx / openpyxl / PyMuPDF / Pillow）<br>"
            "② pandoc（通用文档互转主力引擎）<br><br>"
            "<b>缺少依赖，请重新安装完整包。</b>"
        )
        btn_anyway = dlg.addButton("仍继续（会失败，仅测试）", QMessageBox.ButtonRole.ActionRole)
        btn_cancel = dlg.addButton("取消", QMessageBox.ButtonRole.RejectRole)
        dlg.exec()
        clicked = dlg.clickedButton()
        if clicked is btn_anyway:
            return True
        return False

    def _start_conversion(self):
        if not self._ensure_any_engine():
            return
        tasks = self._build_tasks()
        if not tasks:
            QMessageBox.information(self, "提示", "请先添加要转换的文档。")
            return
        if self.worker and self.worker.isRunning():
            return

        self.convert_btn.setEnabled(False)
        self.stop_btn.setEnabled(True)

        self.worker = DocConversionWorker(
            tasks,
            pandoc_path=self.pandoc_path,
            wkhtmltopdf_path=self.wkhtmltopdf_path,
            native_flags=self.native_flags,
        )
        self.worker.progress_signal.connect(self._on_overall_progress)
        self.worker.task_started_signal.connect(self._on_task_started)
        self.worker.task_finished_signal.connect(self._on_task_finished)
        self.worker.all_done_signal.connect(self._on_all_done)
        self.worker.start()

        output_format = self.format_combo.currentData()
        self.log_signal.emit(
            "info",
            f"📄 文档转换启动：共 {len(tasks)} 个文件 → 输出格式 {str(output_format or '').upper()}"
        )

    def _stop_conversion(self):
        if self.worker and self.worker.isRunning():
            self.worker.cancel()
            self.stop_btn.setEnabled(False)
            self.log_signal.emit("warning", "📄 文档转换已被用户停止")

    # ============================================================
    # 进度回调
    # ============================================================
    def _on_overall_progress(self, current, total):
        pass

    def _on_task_started(self, name):
        self.task_monitor_signal.emit(name)
        self.log_signal.emit("info", f"📄 正在转换: {name}")

    def _on_task_finished(self, result: DocConversionResult):
        basename = os.path.basename(getattr(result.task, "input_path", ""))
        msg = (getattr(result, "message", "") or "").strip()
        first_line = msg.splitlines()[0] if msg else ""
        if len(first_line) > 120:
            first_line = first_line[:117] + "..."
        if getattr(result, "success", False):
            self.log_signal.emit("success", f"📄 {basename} — {first_line or '转换成功'}")
        else:
            self.log_signal.emit("error", f"📄 {basename} — {first_line or '转换失败'}")

    def _on_all_done(self):
        self.convert_btn.setEnabled(True)
        self.stop_btn.setEnabled(False)
        self.task_monitor_signal.emit("")
        self.log_signal.emit("info", "📄 文档转换全部完成")

    # ==================== 主题 / 默认输出目录支持 ====================

    @Slot(str)
    def set_default_output_dir(self, path: str):
        self._default_output_dir = (path or "").strip()
        current = self.output_path_edit.text().strip()
        if not current:
            self.output_path_edit.setText(self._default_output_dir)

    @Slot(dict)
    def reapply_theme(self, colors: dict):
        self.theme_colors = colors
        self._apply_widget_styles()

    def _apply_widget_styles(self):
        c = self.theme_colors

        # --- 标题 ---
        root_layout = self.layout()
        first_vl = root_layout.itemAt(0)
        if first_vl and first_vl.layout():
            tl = first_vl.layout()
            if tl.count() >= 2:
                t_lab = tl.itemAt(0).widget()
                s_lab = tl.itemAt(1).widget()
                if isinstance(t_lab, QLabel):
                    t_lab.setStyleSheet(f"font-size: 22px; font-weight: bold; color: {c['text']};")
                if isinstance(s_lab, QLabel):
                    s_lab.setStyleSheet(f"font-size: 13px; color: {c['text_secondary']};")

        # --- 分组框通用样式 ---
        gb_style = f"""
            QGroupBox {{
                font-weight: bold;
                color: {c['text']};
                border: 1px solid {c['border']};
                border-radius: 8px;
                margin-top: 10px;
                padding-top: 16px;
                background-color: {c['card']};
            }}
            QGroupBox::title {{
                subcontrol-origin: margin;
                left: 12px;
                padding: 0 6px;
            }}
        """
        for gb in self.findChildren(QGroupBox):
            gb.setStyleSheet(gb_style)

        # --- 文件列表 ---
        self.file_list.setStyleSheet(f"""
            QListWidget {{
                background-color: {c['bg']};
                border: 1px solid {c['border']};
                border-radius: 6px;
                color: {c['text']};
                padding: 4px;
            }}
            QListWidget::item {{
                padding: 8px;
                border-bottom: 1px solid {c['border']};
            }}
            QListWidget::item:selected {{
                background-color: {c['primary']};
                color: white;
            }}
        """)

        # --- 文件操作按钮 ---
        btn_file_op = f"""
            QPushButton {{
                background-color: {c['card_hover']};
                color: {c['text']};
                border: 1px solid {c['border']};
                border-radius: 4px;
                padding: 8px 12px;
            }}
            QPushButton:hover {{
                background-color: {c['primary']};
                color: white;
                border-color: {c['primary']};
            }}
        """
        for btn in (self.add_btn, self.remove_btn, self.clear_btn):
            btn.setStyleSheet(btn_file_op)

        # --- 输入控件通用 ---
        input_style = f"""
            QComboBox, QLineEdit, QSpinBox {{
                background-color: {c['bg']};
                border: 1px solid {c['border']};
                border-radius: 4px;
                padding: 6px 10px;
                color: {c['text']};
                min-height: 24px;
            }}
            QComboBox:hover, QLineEdit:hover, QSpinBox:hover,
            QComboBox:focus, QLineEdit:focus, QSpinBox:focus {{
                border: 1px solid {c['primary']};
            }}
            QComboBox::drop-down {{
                border: none;
                width: 24px;
            }}
            QComboBox QAbstractItemView {{
                background-color: {c['card']};
                border: 1px solid {c['border']};
                selection-background-color: {c['primary']};
                color: {c['text']};
            }}
        """
        for w in (self.format_combo, self.output_path_edit, self.dpi_spin):
            try:
                w.setStyleSheet(input_style)
            except Exception:
                pass

        # --- 浏览按钮 ---
        self.output_browse_btn.setStyleSheet(f"""
            QPushButton {{
                background-color: {c['card_hover']};
                color: {c['text']};
                border: 1px solid {c['border']};
                border-radius: 4px;
                padding: 6px 10px;
            }}
            QPushButton:hover {{
                background-color: {c['primary']};
                color: white;
                border-color: {c['primary']};
            }}
        """)

        # --- 底部操作按钮 ---
        self.convert_btn.setStyleSheet(f"""
            QPushButton {{
                background-color: {c["success"]};
                color: white;
                font-weight: bold;
                font-size: 15px;
                border-radius: 8px;
                border: none;
            }}
            QPushButton:hover {{
                background-color: {c.get('success_hover', '#2eb872')};
            }}
            QPushButton:disabled {{
                background-color: {c.get('btn_disabled_bg', '#2a3a4a')};
                color: {c.get('btn_disabled_text', '#555555')};
            }}
        """)
        self.stop_btn.setStyleSheet(f"""
            QPushButton {{
                background-color: {c["accent"]};
                color: white;
                font-weight: bold;
                font-size: 15px;
                border-radius: 8px;
                border: none;
            }}
            QPushButton:hover {{
                background-color: {c.get('accent_hover', '#c0392b')};
            }}
            QPushButton:disabled {{
                background-color: {c.get('btn_disabled_bg', '#2a3a4a')};
                color: {c.get('btn_disabled_text', '#555555')};
            }}
        """)
        self.open_dir_btn.setStyleSheet(f"""
            QPushButton {{
                background-color: {c["card"]};
                color: {c["text"]};
                font-weight: 500;
                font-size: 13px;
                border-radius: 8px;
                border: 1px solid {c["border"]};
            }}
            QPushButton:hover {{
                background-color: {c["card_hover"]};
                border-color: {c["success"]};
            }}
        """)